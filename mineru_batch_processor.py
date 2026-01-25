#!/usr/bin/env python3
"""
MinerU 批量处理器 - 生产级方案
支持：并行处理、文件拆分、结果合并、进度监控
"""
import json
import asyncio
import aiohttp
import random
from pathlib import Path
from typing import List, Dict, Optional, Callable
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor
import time

try:
    from PyPDF2 import PdfReader, PdfWriter
    from pptx import Presentation
    from docx import Document
except ImportError:
    print("❌ 请安装依赖: uv pip install PyPDF2 python-pptx python-docx")
    exit(1)


class FileChunker:
    """文件拆分器"""
    
    MAX_PAGES = 600  # 最大页数限制
    
    @staticmethod
    def split_pdf(file_path: str, output_dir: str) -> List[str]:
        """拆分PDF"""
        reader = PdfReader(file_path)
        total_pages = len(reader.pages)
        
        if total_pages <= FileChunker.MAX_PAGES:
            return [file_path]
        
        chunks = []
        output_path = Path(output_dir)
        output_path.mkdir(exist_ok=True)
        
        file_name = Path(file_path).stem
        chunk_count = (total_pages + FileChunker.MAX_PAGES - 1) // FileChunker.MAX_PAGES
        
        print(f"📄 拆分PDF: {total_pages}页 → {chunk_count}个文件")
        
        for i in range(chunk_count):
            start = i * FileChunker.MAX_PAGES
            end = min((i + 1) * FileChunker.MAX_PAGES, total_pages)
            
            writer = PdfWriter()
            for page_num in range(start, end):
                writer.add_page(reader.pages[page_num])
            
            chunk_path = output_path / f"{file_name}_chunk_{i+1}.pdf"
            with open(chunk_path, 'wb') as f:
                writer.write(f)
            
            chunks.append(str(chunk_path))
            print(f"  ✅ 分片 {i+1}/{chunk_count}: {end-start}页")
        
        return chunks
    
    @staticmethod
    def split_pptx(file_path: str, output_dir: str) -> List[str]:
        """拆分PPTX"""
        prs = Presentation(file_path)
        total_slides = len(prs.slides)
        
        if total_slides <= FileChunker.MAX_PAGES:
            return [file_path]
        
        chunks = []
        output_path = Path(output_dir)
        output_path.mkdir(exist_ok=True)
        
        file_name = Path(file_path).stem
        chunk_count = (total_slides + FileChunker.MAX_PAGES - 1) // FileChunker.MAX_PAGES
        
        print(f"📊 拆分PPTX: {total_slides}页 → {chunk_count}个文件")
        
        for i in range(chunk_count):
            start = i * FileChunker.MAX_PAGES
            end = min((i + 1) * FileChunker.MAX_PAGES, total_slides)
            
            new_prs = Presentation()
            new_prs.slide_width = prs.slide_width
            new_prs.slide_height = prs.slide_height
            
            for slide_num in range(start, end):
                slide = prs.slides[slide_num]
                new_prs.slides.add_slide(slide.slide_layout)
            
            chunk_path = output_path / f"{file_name}_chunk_{i+1}.pptx"
            new_prs.save(str(chunk_path))
            
            chunks.append(str(chunk_path))
            print(f"  ✅ 分片 {i+1}/{chunk_count}: {end-start}页")
        
        return chunks
    
    @staticmethod
    def split_docx(file_path: str, output_dir: str) -> List[str]:
        """拆分DOCX（按段落估算）"""
        doc = Document(file_path)
        total_paragraphs = len(doc.paragraphs)
        
        # 估算：每页约5段落
        estimated_pages = total_paragraphs // 5
        
        if estimated_pages <= FileChunker.MAX_PAGES:
            return [file_path]
        
        chunks = []
        output_path = Path(output_dir)
        output_path.mkdir(exist_ok=True)
        
        file_name = Path(file_path).stem
        paras_per_chunk = FileChunker.MAX_PAGES * 5
        chunk_count = (total_paragraphs + paras_per_chunk - 1) // paras_per_chunk
        
        print(f"📝 拆分DOCX: ~{estimated_pages}页 → {chunk_count}个文件")
        
        for i in range(chunk_count):
            start = i * paras_per_chunk
            end = min((i + 1) * paras_per_chunk, total_paragraphs)
            
            new_doc = Document()
            for para_num in range(start, end):
                new_doc.add_paragraph(doc.paragraphs[para_num].text)
            
            chunk_path = output_path / f"{file_name}_chunk_{i+1}.docx"
            new_doc.save(str(chunk_path))
            
            chunks.append(str(chunk_path))
            print(f"  ✅ 分片 {i+1}/{chunk_count}")
        
        return chunks
    
    @staticmethod
    def split_file(file_path: str, output_dir: str = "./chunks") -> List[str]:
        """自动识别并拆分文件"""
        suffix = Path(file_path).suffix.lower()
        
        if suffix == '.pdf':
            return FileChunker.split_pdf(file_path, output_dir)
        elif suffix in ['.pptx', '.ppt']:
            return FileChunker.split_pptx(file_path, output_dir)
        elif suffix in ['.docx', '.doc']:
            return FileChunker.split_docx(file_path, output_dir)
        else:
            return [file_path]


class ResultMerger:
    """结果合并器 - 处理MinerU API返回的完整结果"""
    
    @staticmethod
    async def download_file(session: aiohttp.ClientSession, url: str, output_path: str):
        """下载文件"""
        async with session.get(url) as resp:
            with open(output_path, 'wb') as f:
                f.write(await resp.read())
    
    @staticmethod
    async def download_and_extract_results(results: List[Dict], output_dir: str) -> List[Dict]:
        """
        下载并解压所有结果
        
        MinerU API返回结构:
        {
            'full_zip_url': 'https://...zip',  # 完整压缩包
            'md_url': 'https://...md',          # Markdown文件
            'md_content_url': 'https://...',    # Markdown内容
            'layout_tree_url': 'https://...'    # 布局树
        }
        """
        import zipfile
        from pathlib import Path
        
        output_path = Path(output_dir)
        output_path.mkdir(exist_ok=True, parents=True)
        
        extracted_results = []
        
        async with aiohttp.ClientSession() as session:
            for i, result in enumerate(results, 1):
                if result['status'] != 'success':
                    continue
                
                data = result['result']
                chunk_dir = output_path / f"chunk_{i}"
                chunk_dir.mkdir(exist_ok=True)
                
                print(f"📥 下载分片 {i}...")
                
                # 下载完整压缩包
                if 'full_zip_url' in data:
                    zip_path = chunk_dir / "result.zip"
                    await ResultMerger.download_file(session, data['full_zip_url'], str(zip_path))
                    
                    # 解压
                    with zipfile.ZipFile(zip_path, 'r') as zip_ref:
                        zip_ref.extractall(chunk_dir)
                    
                    print(f"  ✅ 解压完成: {chunk_dir}")
                
                # 下载Markdown
                md_content = None
                if 'md_content_url' in data:
                    async with session.get(data['md_content_url']) as resp:
                        md_content = await resp.text()
                elif 'md_url' in data:
                    async with session.get(data['md_url']) as resp:
                        md_content = await resp.text()
                
                extracted_results.append({
                    'chunk_id': i,
                    'chunk_dir': str(chunk_dir),
                    'md_content': md_content,
                    'data': data
                })
        
        return extracted_results
    
    @staticmethod
    def merge_markdown_files(extracted_results: List[Dict], output_file: str):
        """合并Markdown内容"""
        with open(output_file, 'w', encoding='utf-8') as f:
            for i, result in enumerate(extracted_results, 1):
                if i > 1:
                    f.write("\n\n" + "="*60 + "\n\n")
                
                f.write(f"# 分片 {i}\n\n")
                
                if result['md_content']:
                    f.write(result['md_content'])
                else:
                    f.write("*内容为空*\n")
        
        print(f"✅ Markdown合并完成: {output_file}")
    
    @staticmethod
    def merge_images(extracted_results: List[Dict], output_dir: str):
        """合并所有图片到统一目录"""
        import shutil
        
        images_dir = Path(output_dir) / "images"
        images_dir.mkdir(exist_ok=True, parents=True)
        
        image_count = 0
        
        for result in extracted_results:
            chunk_dir = Path(result['chunk_dir'])
            
            # 查找所有图片文件
            for img_file in chunk_dir.rglob("*.png"):
                new_name = f"chunk_{result['chunk_id']}_{img_file.name}"
                shutil.copy(img_file, images_dir / new_name)
                image_count += 1
            
            for img_file in chunk_dir.rglob("*.jpg"):
                new_name = f"chunk_{result['chunk_id']}_{img_file.name}"
                shutil.copy(img_file, images_dir / new_name)
                image_count += 1
        
        print(f"✅ 图片合并完成: {image_count} 个文件 → {images_dir}")
    
    @staticmethod
    def merge_json_metadata(extracted_results: List[Dict], output_file: str):
        """合并JSON元数据"""
        merged = {
            'total_chunks': len(extracted_results),
            'merged_at': datetime.now().isoformat(),
            'chunks': [
                {
                    'chunk_id': r['chunk_id'],
                    'chunk_dir': r['chunk_dir'],
                    'has_content': r['md_content'] is not None,
                    'urls': r['data']
                }
                for r in extracted_results
            ]
        }
        
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(merged, f, indent=2, ensure_ascii=False)
        
        print(f"✅ JSON元数据合并完成: {output_file}")


class ProgressMonitor:
    """进度监控器"""
    
    def __init__(self, total: int):
        self.total = total
        self.completed = 0
        self.failed = 0
        self.start_time = time.time()
    
    def update(self, success: bool = True):
        """更新进度"""
        if success:
            self.completed += 1
        else:
            self.failed += 1
        
        self.print_progress()
    
    def print_progress(self):
        """打印进度"""
        elapsed = time.time() - self.start_time
        total_done = self.completed + self.failed
        percent = (total_done / self.total) * 100
        
        print(f"\r⏳ 进度: {total_done}/{self.total} ({percent:.1f}%) | "
              f"✅ {self.completed} | ❌ {self.failed} | "
              f"⏱️  {elapsed:.1f}s", end='', flush=True)
        
        if total_done == self.total:
            print()  # 换行


class MinerUBatchProcessor:
    """MinerU 批量处理器"""
    
    def __init__(self, tokens_file='all_tokens.json', max_workers=10):
        """
        初始化
        
        Args:
            tokens_file: Token文件
            max_workers: 最大并行度
        """
        self.tokens_file = tokens_file
        self.max_workers = max_workers
        self.tokens = self._load_tokens()
        self.base_url = 'https://mineru.net/api/v4'
        
        if not self.tokens:
            raise ValueError("未找到Token，请先运行 batch_login.py")
        
        print(f"✅ 已加载 {len(self.tokens)} 个账户")
        print(f"⚙️  最大并行度: {max_workers}")
    
    def _load_tokens(self) -> Dict:
        """加载Token"""
        try:
            with open(self.tokens_file, 'r') as f:
                return json.load(f)
        except FileNotFoundError:
            return {}
    
    def _get_random_token(self) -> str:
        """随机选择Token"""
        email = random.choice(list(self.tokens.keys()))
        return self.tokens[email]['token']
    
    async def _process_single_file(self, session: aiohttp.ClientSession, 
                                   file_url: str, file_id: str) -> Dict:
        """处理单个文件（异步）"""
        token = self._get_random_token()
        headers = {
            'authorization': f'Bearer {token}',
            'content-type': 'application/json'
        }
        
        # 创建任务
        data = {'url': file_url, 'model_version': 'vlm'}
        
        async with session.post(
            f"{self.base_url}/extract/task",
            headers=headers,
            json=data
        ) as resp:
            result = await resp.json()
            
            if result['code'] != 0:
                return {'file_id': file_id, 'status': 'failed', 'error': result}
            
            task_id = result['data']['task_id']
        
        # 轮询结果
        max_attempts = 60
        for _ in range(max_attempts):
            await asyncio.sleep(5)
            
            async with session.get(
                f"{self.base_url}/extract/task/{task_id}",
                headers=headers
            ) as resp:
                result = await resp.json()
                
                if result['code'] == 0:
                    data = result['data']
                    state = data.get('state')
                    
                    if state == 'done':
                        return {
                            'file_id': file_id,
                            'status': 'success',
                            'task_id': task_id,
                            'result': data
                        }
                    elif state == 'failed':
                        return {
                            'file_id': file_id,
                            'status': 'failed',
                            'error': data.get('err_msg')
                        }
        
        return {'file_id': file_id, 'status': 'timeout'}
    
    async def process_files_async(self, files: List[Dict]) -> List[Dict]:
        """异步批量处理文件"""
        monitor = ProgressMonitor(len(files))
        
        async with aiohttp.ClientSession() as session:
            semaphore = asyncio.Semaphore(self.max_workers)
            
            async def process_with_semaphore(file_info):
                async with semaphore:
                    result = await self._process_single_file(
                        session,
                        file_info['url'],
                        file_info['id']
                    )
                    monitor.update(result['status'] == 'success')
                    return result
            
            tasks = [process_with_semaphore(f) for f in files]
            results = await asyncio.gather(*tasks)
        
        return results
    
    def process_files(self, files: List[Dict]) -> List[Dict]:
        """批量处理文件（同步入口）"""
        print(f"\n🚀 开始处理 {len(files)} 个文件")
        print(f"⚙️  并行度: {self.max_workers}")
        print()
        
        results = asyncio.run(self.process_files_async(files))
        
        # 统计
        success = sum(1 for r in results if r['status'] == 'success')
        failed = sum(1 for r in results if r['status'] == 'failed')
        timeout = sum(1 for r in results if r['status'] == 'timeout')
        
        print(f"\n📊 处理完成:")
        print(f"  ✅ 成功: {success}")
        print(f"  ❌ 失败: {failed}")
        print(f"  ⏱️  超时: {timeout}")
        
        return results
    
    def process_large_file(self, file_path: str, output_dir: str = "./output") -> Dict:
        """
        处理大文件（自动拆分、并行处理、完整合并）
        
        Args:
            file_path: 文件路径
            output_dir: 输出目录
        
        Returns:
            处理结果
        """
        print(f"\n📄 处理大文件: {file_path}")
        
        # 1. 拆分文件
        chunks = FileChunker.split_file(file_path, "./chunks")
        print(f"📦 拆分完成: {len(chunks)} 个分片")
        
        if len(chunks) == 1:
            print("💡 文件无需拆分，直接处理")
        
        # 2. 上传分片（这里需要实际上传逻辑）
        # TODO: 实现文件上传到CDN
        print("\n⚠️  注意: 需要先上传分片到CDN，获取URL")
        print("💡 提示: 使用 upload_chunks() 方法上传")
        
        files = [
            {'id': f'chunk_{i}', 'url': f'https://example.com/{Path(c).name}'}
            for i, c in enumerate(chunks)
        ]
        
        # 3. 并行处理
        results = self.process_files(files)
        
        # 4. 下载并解压所有结果
        print(f"\n📥 下载并解压结果...")
        extracted_results = asyncio.run(
            ResultMerger.download_and_extract_results(results, output_dir)
        )
        
        # 5. 合并所有内容
        output_path = Path(output_dir)
        file_name = Path(file_path).stem
        
        print(f"\n🔗 合并结果...")
        
        # 合并Markdown
        md_file = output_path / f"{file_name}_merged.md"
        ResultMerger.merge_markdown_files(extracted_results, str(md_file))
        
        # 合并图片
        ResultMerger.merge_images(extracted_results, output_dir)
        
        # 合并元数据
        json_file = output_path / f"{file_name}_metadata.json"
        ResultMerger.merge_json_metadata(extracted_results, str(json_file))
        
        return {
            'total_chunks': len(chunks),
            'success': len(extracted_results),
            'failed': len(results) - len(extracted_results),
            'output_files': {
                'markdown': str(md_file),
                'images': str(output_path / "images"),
                'metadata': str(json_file)
            }
        }


# 使用示例
if __name__ == '__main__':
    import sys
    
    if len(sys.argv) < 2:
        print("用法:")
        print("  批量处理: python3 mineru_batch_processor.py batch <file1> <file2> ...")
        print("  大文件处理: python3 mineru_batch_processor.py large <file>")
        sys.exit(1)
    
    mode = sys.argv[1]
    processor = MinerUBatchProcessor(max_workers=10)
    
    if mode == 'batch':
        # 批量处理多个文件
        files = [
            {'id': f'file_{i}', 'url': f'https://example.com/{Path(f).name}'}
            for i, f in enumerate(sys.argv[2:])
        ]
        processor.process_files(files)
    
    elif mode == 'large':
        # 处理大文件
        file_path = sys.argv[2]
        result = processor.process_large_file(file_path)
        print(f"\n✅ 处理完成: {result}")
